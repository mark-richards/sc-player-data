"""
create_presentation.py

Generates the SuperCoach analytics work presentation as a .pptx file.
Uses Reece Trade brand guidelines: colours, fonts, and slide layouts.

Run:
    pip install python-pptx
    py create_presentation.py

Output: supercoach_presentation.pptx (in the project root)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ---------------------------------------------------------------------------
# Brand colours
# ---------------------------------------------------------------------------
REECE_BLUE   = RGBColor(0x00, 0x30, 0x57)   # #003057
TRUE_BLUE    = RGBColor(0x6F, 0xB1, 0xC8)   # #6FB1C8
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)   # #FFFFFF
CHROME_GREY  = RGBColor(0xED, 0xED, 0xED)   # #EDEDED
TRADIE_GREY  = RGBColor(0x57, 0x57, 0x56)   # #575756

# Slide dimensions — standard widescreen 16:9
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def set_bg(slide, colour: RGBColor):
    """Fill slide background with a solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = colour


def add_textbox(slide, text, left, top, width, height,
                font_size=Pt(24), bold=False, colour=WHITE,
                align=PP_ALIGN.LEFT, word_wrap=True):
    """Add a simple text box to a slide and return the text frame."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = colour
    run.font.name = "Arial"
    return tf


def add_bullet_textbox(slide, bullets, left, top, width, height,
                       font_size=Pt(18), colour=REECE_BLUE):
    """Add a text box with multiple bullet lines."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for bullet in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = bullet
        run.font.size = font_size
        run.font.color.rgb = colour
        run.font.name = "Arial"
    return tf


def add_notes(slide, notes_text: str):
    """Populate the speaker notes pane."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


def divider_slide(prs, headline, subheadline="", notes=""):
    """
    Section divider: full Reece Blue background, massive bold white headline,
    optional True Blue subheadline underneath.
    """
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    set_bg(slide, REECE_BLUE)

    # Headline
    add_textbox(
        slide, headline,
        left=Inches(0.6), top=Inches(2.2),
        width=Inches(12.0), height=Inches(2.5),
        font_size=Pt(48), bold=True, colour=WHITE,
        align=PP_ALIGN.LEFT,
    )

    # Subheadline
    if subheadline:
        add_textbox(
            slide, subheadline,
            left=Inches(0.6), top=Inches(4.6),
            width=Inches(12.0), height=Inches(0.8),
            font_size=Pt(22), bold=True, colour=TRUE_BLUE,
            align=PP_ALIGN.LEFT,
        )

    if notes:
        add_notes(slide, notes)

    return slide


def content_slide(prs, heading, subheading, bullets, notes="", bg=WHITE):
    """
    Content slide: white/Chrome Grey background, Reece Blue heading,
    True Blue subheading, bullet body.
    """
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    set_bg(slide, bg)

    # Heading
    add_textbox(
        slide, heading,
        left=Inches(0.5), top=Inches(0.35),
        width=Inches(12.3), height=Inches(0.9),
        font_size=Pt(32), bold=True, colour=REECE_BLUE,
        align=PP_ALIGN.LEFT,
    )

    # Subheading
    if subheading:
        add_textbox(
            slide, subheading,
            left=Inches(0.5), top=Inches(1.2),
            width=Inches(12.3), height=Inches(0.5),
            font_size=Pt(18), bold=True, colour=TRUE_BLUE,
            align=PP_ALIGN.LEFT,
        )

    # Bullets
    if bullets:
        bullet_top = Inches(1.85) if subheading else Inches(1.3)
        add_bullet_textbox(
            slide, bullets,
            left=Inches(0.5), top=bullet_top,
            width=Inches(12.3), height=Inches(5.0),
            font_size=Pt(20), colour=REECE_BLUE,
        )

    if notes:
        add_notes(slide, notes)

    return slide


# ---------------------------------------------------------------------------
# Build presentation
# ---------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # ------------------------------------------------------------------
    # Slide 1 — Hook (DIVIDER)
    # ------------------------------------------------------------------
    divider_slide(
        prs,
        headline="Billy Beane had a spreadsheet. I have a machine learning model.",
        subheadline="A personal project. A footy obsession. A lot of Python.",
        notes=(
            "So you've watched Moneyball. Billy Beane walked into a room full of scouts "
            "who trusted their gut, and he said: 'No. We're using the data.' That's kind "
            "of what I've been doing in my living room on weekends — except instead of the "
            "Oakland A's, it's my SuperCoach fantasy footy team. And instead of on-base "
            "percentage, it's a LightGBM model predicting AFL player scores. Let me show "
            "you what I built."
        ),
    )

    # ------------------------------------------------------------------
    # Slide 2 — Show of hands (CONTENT)
    # ------------------------------------------------------------------
    content_slide(
        prs,
        heading="Let's see where we're at.",
        subheading="Show of hands.",
        bullets=[
            "🙋 Who here plays fantasy sport? (SuperCoach, AFL Fantasy, NRL, anything)",
            "🙋 Who's done a personal coding or data project outside of work?",
            "🙋 Who thought Moneyball was actually a great film?",
        ],
        notes=(
            "Before I dive in — let's get a feel for the room. Hands up if you play any "
            "kind of fantasy sport. Keep it up if you've ever built something in your own "
            "time — a script, a dashboard, a little side project. And of course — who "
            "genuinely enjoyed Moneyball? Good. You'll get the references today."
        ),
    )

    # ------------------------------------------------------------------
    # Slide 3 — What is SuperCoach? (CONTENT)
    # ------------------------------------------------------------------
    content_slide(
        prs,
        heading="SuperCoach is fantasy footy. Here's how it works.",
        subheading="800+ AFL players. One job: pick the best team.",
        bullets=[
            "Each real AFL player earns points based on their stats each week (kicks, marks, tackles)",
            "You pick a squad of 23 — on-field starters + bench",
            "Trade players in and out each week to maximise your score",
            "Compete against others in a league — highest score wins",
        ],
        notes=(
            "For the uninitiated — SuperCoach is a fantasy footy game. You draft or pick a "
            "squad of AFL players, and each week they earn points based on what they actually "
            "do on the field. Kicks, marks, tackles, handballs — all converted into a score. "
            "The goal is to have the highest-scoring team in your league each week. There are "
            "about 800 players to choose from. The question is: who do you pick?"
        ),
    )

    # ------------------------------------------------------------------
    # Slide 4 — The problem (CONTENT)
    # ------------------------------------------------------------------
    content_slide(
        prs,
        heading="I was losing. And I was guessing.",
        subheading="There had to be a better way.",
        bullets=[
            "Every week: who's underpriced? Who's about to explode?",
            "800 players — impossible to track manually",
            "Everyone else was using gut feel and Twitter hot takes",
            "I wanted an actual edge",
        ],
        notes=(
            "Here's the honest truth — I was making decisions based on vibes. I'd scroll "
            "Twitter, read a footy forum, and pick whoever sounded hot. Sound familiar? "
            "That's exactly what the scouts in Moneyball were doing. And just like Billy "
            "Beane, I got sick of it. I figured: I work in data. I know how to build stuff. "
            "Why am I not using that here?"
        ),
    )

    # ------------------------------------------------------------------
    # Slide 5 — What I built: overview (CONTENT)
    # ------------------------------------------------------------------
    content_slide(
        prs,
        heading="I built a full analytics pipeline. From scratch.",
        subheading="Scrape → process → predict → publish.",
        bullets=[
            "🔍  Scrape   — FanFooty match data + SuperCoach API (every week, automated)",
            "⚙️   Process  — clean, merge, engineer 76 features per player per round",
            "🧠  Predict  — LightGBM ML model → player score forecasts",
            "📬  Publish  — weekly newsletter + live web dashboard",
        ],
        notes=(
            "So here's what I actually built. It starts with a web scraper that pulls match "
            "data every week. That feeds into a data pipeline that cleans and processes "
            "everything. The clean data goes into a machine learning model that predicts "
            "player scores. And the output? A weekly newsletter and a live web dashboard I "
            "built myself. All in Python, all running on my own machine, all in my own time."
        ),
    )

    # ------------------------------------------------------------------
    # Slide 6 — Deep dive: data pipeline (CONTENT)
    # ------------------------------------------------------------------
    content_slide(
        prs,
        heading="Every week, the machine wakes up.",
        subheading="Automated. Repeatable. No manual work.",
        bullets=[
            "Scrapes FanFooty.com.au — match stats for every AFL game, back to 2011",
            "Hits the SuperCoach API — live player prices, ownership, league data",
            "Cleans and merges ~15 years of player data into one master dataset",
            "Engineers 76 features per player: form, ceiling rate, time on ground, and more",
            "💬 MONEYBALL Q: What stat did they use instead of batting average? → On-base %.\n    SuperCoach equivalent: Ceiling Index — how often a player scores 120+.",
        ],
        notes=(
            "The data pipeline is the engine room. Every week it scrapes two sources — "
            "FanFooty for match stats going back to 2011, and the SuperCoach API for live "
            "league data. That all gets cleaned and merged into one master dataset. Then I "
            "engineer 76 features per player — things like rolling form, how often they score "
            "above 120, their time on ground trends. That's what the model actually learns from.\n\n"
            "MONEYBALL QUESTION (pause here): In Moneyball, what was the stat they used "
            "instead of batting average? On-base percentage. In SuperCoach terms, that's the "
            "Ceiling Index — how often a player scores above 120 points. It's not about "
            "averages. It's about ceiling."
        ),
    )

    # ------------------------------------------------------------------
    # Slide 7 — Deep dive: ML model (CONTENT)
    # ------------------------------------------------------------------
    content_slide(
        prs,
        heading="The model knows more than your gut does.",
        subheading="LightGBM. 76 features. Trained on 15 years of data.",
        bullets=[
            "Model: LightGBM — gradient boosting (same family as Kaggle competition winners)",
            "Trained on AFL data 2011–2025 using an expanding window",
            "Predicts each player's SuperCoach score for the upcoming round",
            "Error rate: ~10.4 points on average — within one good kick of the truth",
            "❓ TRIVIA: How many players are in the SC pool? (pause for guesses) → ~800",
        ],
        notes=(
            "The model is a LightGBM — a gradient boosting algorithm that's excellent on "
            "tabular data. I trained it on 15 years of AFL match data using an expanding "
            "window — each year it learns from all previous years, then predicts the next. "
            "The average prediction error is about 10 points per player per round. That's "
            "the equivalent of about one decent kick. It's not perfect, but it's a lot "
            "better than a gut feel.\n\n"
            "TRIVIA (pause): How many players are in the SuperCoach player pool? "
            "Have a guess. (pause) About 800. And my model has a prediction for every "
            "single one of them, every single week."
        ),
    )

    # ------------------------------------------------------------------
    # Slide 8 — Demo cue (DIVIDER)
    # ------------------------------------------------------------------
    divider_slide(
        prs,
        headline="Let me just show you.",
        subheadline="Live demo — 2 minutes.",
        notes=(
            "DEMO CUE — switch to browser now.\n\n"
            "Show in order:\n"
            "1. Web dashboard — live league standings, player scores, team view.\n"
            "   Say: 'This is running on Flask, a Python web framework. Auto-refreshes "
            "from the API. I built this myself.'\n\n"
            "2. Weekly newsletter — open latest output.\n"
            "   Say: 'Every week this generates automatically. It recommends which free "
            "agents to pick up, flags players to drop, and highlights the ceiling chasers "
            "— the players most likely to go huge.'\n\n"
            "Return to slides after ~2 minutes."
        ),
    )

    # ------------------------------------------------------------------
    # Slide 9 — Results + quiz (CONTENT)
    # ------------------------------------------------------------------
    content_slide(
        prs,
        heading="Does it actually work? Let's see.",
        subheading="You tell me.",
        bullets=[
            "❓ QUIZ: What % of SC scores break 120 in a given round? (pause for guesses)",
            "→ ~15% of scores in any round hit the 120+ ceiling threshold",
            "Model's top picks hit ceiling at roughly 2× the average rate",
            "Prediction error: 10.4 MAE vs 12.8 baseline — an 18.9% improvement",
            "Better decisions. Every week.",
        ],
        notes=(
            "QUIZ (show question first, pause for guesses, then reveal bullets):\n"
            "What percentage of SuperCoach scores do you think are above 120 points "
            "in a given round?\n\n"
            "So does it actually work? I think so. About 15% of SC scores break 120 "
            "in any given round — those are the ceiling scorers you want in your team. "
            "My model's top picks hit that rate at roughly double the baseline. And the "
            "overall prediction error dropped from 12.8 to 10.4 — that's an 18.9% "
            "improvement over the naive approach.\n\n"
            "NOTE: Verify the 15% and 2x figures against your actual data before presenting."
        ),
        bg=CHROME_GREY,
    )

    # ------------------------------------------------------------------
    # Slide 10a — Moneyball transition (DIVIDER)
    # ------------------------------------------------------------------
    divider_slide(
        prs,
        headline="This is Moneyball. But for footy.",
        notes="Let this land for a moment. Then click to the next slide.",
    )

    # ------------------------------------------------------------------
    # Slide 10b — Moneyball parallels (CONTENT)
    # ------------------------------------------------------------------
    content_slide(
        prs,
        heading="The parallels are pretty clear.",
        subheading="Data beats gut feel. Every time.",
        bullets=[
            "Moneyball: scouts used gut feel  →  SC: coaches use Twitter hot takes",
            "Moneyball: on-base % (OBP)       →  SC: Ceiling Index (CI%)",
            "Moneyball: undervalued players   →  SC: underpriced free agents",
            "Moneyball: systematic selection  →  SC: weekly newsletter picks",
            "💬 MONEYBALL Q: Where in your day job are you still relying on gut feel?",
        ],
        notes=(
            "Here's the Moneyball moment. In the film, Billy Beane didn't just use better "
            "stats — he built a system that found value where others couldn't see it. "
            "That's exactly what this does. Instead of on-base percentage, I use Ceiling "
            "Index. Instead of undervalued players, I'm looking for underpriced free agents. "
            "The principle is identical: replace intuition with a repeatable, data-driven "
            "process.\n\n"
            "QUESTION (pause): Pete Brand in Moneyball used first principles thinking to "
            "see what the scouts couldn't. Where in your day job are you still relying on "
            "gut feel when you could be using data?"
        ),
    )

    # ------------------------------------------------------------------
    # Slide 11 — Skills transfer (CONTENT)
    # ------------------------------------------------------------------
    content_slide(
        prs,
        heading="Every skill here maps to your work.",
        subheading="This isn't just a footy project.",
        bullets=[
            "Web scraping       →  API integrations, data extraction tasks",
            "Data pipelines     →  Power Automate flows, ETL, Power BI refreshes",
            "Machine learning   →  forecasting, classification, anomaly detection",
            "Flask dashboard    →  internal reporting tools, lightweight web apps",
            "Automated reports  →  any scheduled alert or report you've ever built",
        ],
        notes=(
            "None of this is exotic. Web scraping is just API calls and HTML parsing — "
            "you do that every time you pull data into Power BI. The data pipeline is the "
            "same logic as any Power Automate flow. The ML model is the same class of tool "
            "we use for forecasting at work. The newsletter is a scheduled report. The "
            "skills are identical — I just learned them by doing something I actually "
            "cared about."
        ),
        bg=CHROME_GREY,
    )

    # ------------------------------------------------------------------
    # Slide 12a — Call to action (DIVIDER)
    # ------------------------------------------------------------------
    divider_slide(
        prs,
        headline="Your turn.",
        subheadline="Pick something you care about. Build something small. See where it goes.",
        notes="Let this land. Pause. Then click to the next slide.",
    )

    # ------------------------------------------------------------------
    # Slide 12b — How to start (CONTENT)
    # ------------------------------------------------------------------
    content_slide(
        prs,
        heading="How to start. Seriously.",
        subheading="Start small. Stay curious.",
        bullets=[
            "Pick a problem you actually care about — sport, finance, music, travel, anything",
            "Build the scraper or the CSV first. One step at a time.",
            "Use Claude or ChatGPT to help you write the code — you don't need to know everything",
            "Share it. Even half-finished. That's how you learn.",
            "🙋 Who here has a project idea they've been sitting on?",
        ],
        notes=(
            "I'm not standing up here to show off. I'm standing up here because I want "
            "you to do this too. The biggest thing I learned? Just start. The scraper was "
            "30 lines of Python. The first version of the model was terrible. But I built "
            "on it week by week because I cared about the outcome. Pick something you'd "
            "actually use, and go. And use AI to help you — that's the whole point of the "
            "tools we have now.\n\n"
            "SHOW OF HANDS: Who here has a project idea they've been sitting on? "
            "Keep your hand up if the only thing stopping you is not knowing where to start."
        ),
    )

    # ------------------------------------------------------------------
    # Slide 13 — Q&A closer (DIVIDER)
    # ------------------------------------------------------------------
    divider_slide(
        prs,
        headline="Questions? Let's talk.",
        subheadline="Happy to go deeper — the code, the model, or the footy tips.",
        notes=(
            "That's me done. Ten minutes, one machine learning model, and hopefully at "
            "least one of you is now thinking about a side project. I'm happy to chat "
            "about any of it — the code is on my laptop if you want to see it, and yes, "
            "I will share the newsletter picks if you're brave enough to ask. Over to you."
        ),
    )

    # ------------------------------------------------------------------
    # Save
    # ------------------------------------------------------------------
    output_path = "supercoach_presentation.pptx"
    prs.save(output_path)
    print(f"Saved: {output_path}")
    print(f"    Slides: {len(prs.slides)}")
    print()
    print("Next steps:")
    print("  1. Open in PowerPoint")
    print("  2. Add your images/visuals to each slide")
    print("  3. Replace [X] in slide 2 notes with your team size")
    print("  4. Verify the ~15% and 2x stats on slide 9 against your real data")
    print("  5. Add Reece logo (white keyline) to bottom-right of divider slides")


if __name__ == "__main__":
    build()
